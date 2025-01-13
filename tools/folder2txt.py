import sys
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QVBoxLayout, QPushButton, QFileDialog, QCheckBox, QLabel, QWidget, QMessageBox
)
from pathlib import Path
import os
from PyPDF2 import PdfReader
import docx
from pptx import Presentation


class TextExtractorApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.init_ui()

    def init_ui(self):
        self.setWindowTitle("Text Extractor")
        self.setGeometry(100, 100, 400, 200)

        # Main layout
        layout = QVBoxLayout()

        # Folder selection buttons and labels
        self.source_folder_label = QLabel("Source Folder: Not Selected")
        self.target_folder_label = QLabel("Target Folder: Not Selected")
        layout.addWidget(self.source_folder_label)
        layout.addWidget(self.target_folder_label)

        self.select_source_button = QPushButton("Select Source Folder")
        self.select_source_button.clicked.connect(self.select_source_folder)
        layout.addWidget(self.select_source_button)

        self.select_target_button = QPushButton("Select Target Folder")
        self.select_target_button.clicked.connect(self.select_target_folder)
        layout.addWidget(self.select_target_button)

        # Subfolder checkbox
        self.subfolder_checkbox = QCheckBox("Include Subfolders")
        layout.addWidget(self.subfolder_checkbox)

        # Start button
        self.start_button = QPushButton("Start Extraction")
        self.start_button.clicked.connect(self.start_extraction)
        layout.addWidget(self.start_button)

        # Set central widget
        container = QWidget()
        container.setLayout(layout)
        self.setCentralWidget(container)

        # Variables to store folder paths
        self.source_folder = None
        self.target_folder = None

    def select_source_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "Select Source Folder")
        if folder:
            self.source_folder = Path(folder)
            self.source_folder_label.setText(f"Source Folder: {folder}")

    def select_target_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "Select Target Folder")
        if folder:
            self.target_folder = Path(folder)
            self.target_folder_label.setText(f"Target Folder: {folder}")

    def start_extraction(self):
        if not self.source_folder or not self.target_folder:
            QMessageBox.warning(self, "Error", "Please select both source and target folders.")
            return

        include_subfolders = self.subfolder_checkbox.isChecked()
        self.extract_text_from_folder(self.source_folder, self.target_folder, include_subfolders)
        QMessageBox.information(self, "Success", "Text extraction completed!")

    def extract_text_from_folder(self, source_folder, target_folder, include_subfolders):
        # Define file extensions to process
        text_extensions = {".txt", ".md"}
        pdf_extensions = {".pdf"}
        docx_extensions = {".docx"}
        pptx_extensions = {".pptx"}

        # Walk through the folder
        for root, dirs, files in os.walk(source_folder):
            for file in files:
                file_path = Path(root) / file
                file_extension = file_path.suffix.lower()

                # Extract text based on file type
                extracted_text = None
                if file_extension in text_extensions:
                    extracted_text = self.extract_text_from_txt(file_path)
                elif file_extension in pdf_extensions:
                    extracted_text = self.extract_text_from_pdf(file_path)
                elif file_extension in docx_extensions:
                    extracted_text = self.extract_text_from_docx(file_path)
                elif file_extension in pptx_extensions:
                    extracted_text = self.extract_text_from_pptx(file_path)

                # Write extracted text to target folder
                if extracted_text:
                    relative_path = file_path.relative_to(source_folder)
                    target_file_path = target_folder / relative_path.with_suffix(".txt")
                    target_file_path.parent.mkdir(parents=True, exist_ok=True)
                    with open(target_file_path, "w", encoding="utf-8") as f:
                        f.write(extracted_text)

            # If subfolders are not included, break after the first directory
            if not include_subfolders:
                break

    def extract_text_from_txt(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def extract_text_from_pdf(self, file_path):
        try:
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
            return None

    def extract_text_from_docx(self, file_path):
        try:
            doc = docx.Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)
        except Exception as e:
            print(f"Error reading DOCX {file_path}: {e}")
            return None

    def extract_text_from_pptx(self, file_path):
        try:
            presentation = Presentation(file_path)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
            return None


if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = TextExtractorApp()
    window.show()
    sys.exit(app.exec_())
