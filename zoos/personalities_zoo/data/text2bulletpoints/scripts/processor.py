from lollms.types import MSG_OPERATION_TYPE
from lollms.config import TypedConfig, BaseConfig, ConfigTemplate, InstallOption
from lollms.types import MSG_OPERATION_TYPE
from lollms.personality import APScript, AIPersonality
from lollms.helpers import ASCIIColors
from lollms.prompting import LollmsContextDetails

import json
from pathlib import Path
from typing import Callable, Any

class Text2Paragraphs:
    def __init__(self, database_path=None, max_chunk_size=2000):
        self.paragraphs = []
        self.database_path = database_path
        self.max_chunk_size = max_chunk_size
        if database_path is not None:
            if Path(database_path).exists():
                self.load_from_json()

    def chunk_text(self, text):
        paragraphs = text.split("\n")  # Split text by double newlines to separate paragraphs
        if self.max_chunk_size is not None:
            current_chunk = []
            current_chunk_size = 0
            for i,paragraph in enumerate(paragraphs):
                ASCIIColors.yellow(f"Processing paragraph:{i}/{len(paragraphs)}",end="\r")
                if current_chunk_size + len(paragraph) <= self.max_chunk_size:
                    current_chunk.append(paragraph)
                    current_chunk_size += len(paragraph)
                else:
                    self.paragraphs.append("\n".join(current_chunk))
                    current_chunk = [paragraph]
                    current_chunk_size = len(paragraph)
            if current_chunk:
                self.paragraphs.append("\n".join(current_chunk))
        else:
            self.paragraphs.extend(paragraphs)  # Add new paragraphs to the existing ones

    def load_from_json(self, filename=None):
        if filename is None:
            filename = self.database_path
        with open(filename, "r") as file:
            data = json.load(file)
            self.paragraphs = data["paragraphs"]

    def save_to_json(self, filename=None):
        if filename is None:
            filename = self.database_path
        data = {"paragraphs": self.paragraphs}
        with open(filename, "w") as file:
            json.dump(data, file)


class Processor(APScript):
    """
    A class that processes model inputs and outputs.

    Inherits from APScript.
    """

    def __init__(
                 self, 
                 personality: AIPersonality,
                 callback = None,
                ) -> None:
        
        self.word_callback = None
        
        personality_config_template = ConfigTemplate(
            [
                {"name":"database_path","type":"str","value":f"{personality.name}_db.json", "help":"Path to the database"},
                {"name":"max_chunk_size","type":"int","value":512*3, "min":10, "max":personality.config["ctx_size"],"help":"Maximum size of text chunks to vectorize"},
                
                {"name":"max_answer_size","type":"int","value":512, "min":10, "max":personality.config["ctx_size"],"help":"Maximum number of tokens to allow the generator to generate as an answer to your question"},
                
            ]
            )
        personality_config_vals = BaseConfig.from_template(personality_config_template)

        personality_config = TypedConfig(
            personality_config_template,
            personality_config_vals
        )
        super().__init__(
                            personality,
                            personality_config,
                            callback=callback
                        )
        
        self.state = 0

        self.text_store = Text2Paragraphs(
                                    self.personality.lollms_paths.personal_data_path/self.personality_config.database_path,
                                    max_chunk_size=self.personality_config.max_chunk_size
                                    )
        
    @staticmethod        
    def read_pdf_file(file_path):
        import PyPDF2
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text

    @staticmethod
    def read_docx_file(file_path):
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path):
        with open(file_path, 'r') as file:
            data = json.load(file)
        return data
    
    @staticmethod
    def read_csv_file(file_path):
        import csv
        with open(file_path, 'r') as file:
            csv_reader = csv.reader(file)
            lines = [row for row in csv_reader]
        return lines    

    @staticmethod
    def read_html_file(file_path):
        from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    @staticmethod
    def read_pptx_file(file_path):
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content
    
    def build_db(self):
        ASCIIColors.info("-> Vectorizing the database")
        for file in self.personality.text_files:
            try:
                if Path(file).suffix==".pdf":
                    text =  Processor.read_pdf_file(file)
                elif Path(file).suffix==".docx":
                    text =  Processor.read_docx_file(file)
                elif Path(file).suffix==".docx":
                    text =  Processor.read_pptx_file(file)
                elif Path(file).suffix==".json":
                    text =  Processor.read_json_file(file)
                elif Path(file).suffix==".csv":
                    text =  Processor.read_csv_file(file)
                elif Path(file).suffix==".html":
                    text =  Processor.read_html_file(file)
                else:
                    text =  Processor.read_text_file(file)
                    
                self.text_store.chunk_text(text)       
                self.text_store.save_to_json()         
                print(ASCIIColors.color_reset)
                ASCIIColors.success(f"File {file} vectorized successfully")
            except Exception as ex:
                ASCIIColors.error(f"Couldn't vectorize {file}: The vectorizer threw this exception:{ex}")
        ASCIIColors.info("-> Done Vectorizing the database")

    def add_file(self, path, client, callback=None):
        if callback is None and self.callback is not None:
            callback = self.callback
        super().add_file(path, client, callback)
        try:
            self.build_db()
            self.info("File added successfully", callback=callback)
            return True
        except Exception as ex:
            ASCIIColors.error(f"Couldn't vectorize the database: The vectgorizer threw this exception: {ex}")
            return False        

    from lollms.client_session import Client
    def run_workflow(self,  context_details:LollmsContextDetails=None, client:Client=None,  callback: Callable[[str | list | None, MSG_OPERATION_TYPE, str, AIPersonality| None], bool]=None):
        """
        This function generates code based on the given parameters.

        Args:
            context_details (dict): A dictionary containing the following context details for code generation:
                - conditionning (str): The conditioning information.
                - documentation (str): The documentation information.
                - knowledge (str): The knowledge information.
                - user_description (str): The user description information.
                - discussion_messages (str): The discussion messages information.
                - positive_boost (str): The positive boost information.
                - negative_boost (str): The negative boost information.
                - current_language (str): The force language information.
                - fun_mode (str): The fun mode conditionning text
                - ai_prefix (str): The AI prefix information.
            client_id: The client ID for code generation.
            callback (function, optional): The callback function for code generation.

        Returns:
            None
        """
        prompt = context_details.prompt
        previous_discussion_text = context_details.discussion_messages
        # State machine
        self.word_callback = callback
        output =""
        if prompt.strip().lower()=="send_file":
            self.state = 1
            print("Please provide the file name")
            if callback is not None:
                callback("Please provide the file path", MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_SET_CONTENT)
            output = "Please provide the file name"
        elif prompt.strip().lower()=="help":
            ASCIIColors.info("Showing help")
            if callback:
                callback(self.personality.help,MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_SET_CONTENT)
                ASCIIColors.info(help)
            self.state = 0   
        elif prompt.strip().lower()=="show_database":
            if callback:
                callback("Current database\n",MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_ADD_CHUNK)
                print("Current database\n")
            for chunk in self.text_store.paragraphs:
                if callback:
                    callback(chunk+"\n",MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_ADD_CHUNK)
                    print(chunk)
            
            self.state = 0
        elif prompt.strip().lower()=="clear_database":
            database_fill_path:Path = self.personality.lollms_paths.personal_data_path/self.personality_config["database_path"]
            if database_fill_path.exists():
                database_fill_path.unlink()
                self.personality_config.database_path = prompt
                self.personality_config.save()
                self.text_store = Text2Paragraphs(
                                self.personality.lollms_paths.personal_data_path/self.personality_config.database_path,
                                max_chunk_size=self.personality_config.max_chunk_size
                                )   
                if callback is not None:
                    callback("Database file cleared successfully", MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_SET_CONTENT)        
            else:
                if callback is not None:
                    callback("The database file does not exist yet, so you can't clear it", MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_SET_CONTENT)        
            self.state = 0
        elif prompt.strip().lower()=="set_database":
            print("Please provide the database file name")
            if callback is not None:
                callback("Please provide the database file path", MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_SET_CONTENT)
            output = "Please provide the database file name"
            self.state = 2
        elif prompt.strip().lower()=="convert":
            if callback is not None:
                callback("# Full bullet points summary:\n", MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_ADD_CHUNK)
            for i,chunk in enumerate(self.text_store.paragraphs):
                if len(chunk.split())<50:
                    print(chunk)
                    continue
                docs = "\n".join([
                  f"{self.config.start_header_id_template}{self.config.system_message_template}:",
                  "Summarize the following paragraph in the form of bullet points.",
                  "Be concise and only keep most important ideas.",
                  "Use short sentences",
                  f"Paragraph:{chunk}",
                  "Bullet points:\n-"
                ])
                ASCIIColors.error("\n-------------- Documentation -----------------------")
                ASCIIColors.error(docs)
                ASCIIColors.error("----------------------------------------------------")
                print("Thinking")
                output = "-"+self.generate(docs, self.personality_config.max_answer_size)
                self.text_store.paragraphs[i]= output
                if callback is not None:
                    callback(output, MSG_OPERATION_TYPE.MSG_OPERATION_TYPE_ADD_CHUNK)
                    
        elif prompt.strip().lower()=="clear_database":
            (self.personality.lollms_paths.personal_data_path/self.personality_config["database_path"]).unlink()
            self.text_store = Text2Paragraphs(
                            self.personality.lollms_paths.personal_data_path/self.personality_config.database_path,
                            max_chunk_size=self.personality_config.max_chunk_size
                            )            
            self.state = 0
        else:
            if self.state ==1:
                try:
                    pass#self.add_file(prompt, client) # TODO: FIX
                except Exception as ex:
                    ASCIIColors.error(f"Exception: {ex}")
                    output = str(ex)
                self.state=0
            elif self.state ==2:
                try:
                    new_discussion_db_name = Path(prompt)
                    if new_discussion_db_name.exists():
                        self.personality_config.database_path = prompt
                        self.text_store = Text2Paragraphs(
                                        self.personality.lollms_paths.personal_data_path/self.personality_config.database_path,
                                        max_chunk_size=self.personality_config.max_chunk_size
                                        )   
                        
                        self.personality_config.config.save_config()
                    else:
                        output = "Database file not found.\nGoing back to default state."
                except Exception as ex:
                    ASCIIColors.error(f"Exception: {ex}")
                    output = str(ex)
                self.state=0
            else:
                output = "unknown command"
        return output



