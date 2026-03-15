from lollms.types import MSG_OPERATION_TYPE
from lollms.helpers import ASCIIColors
from lollms.utilities import PackageManager
from lollms.config import TypedConfig, BaseConfig, ConfigTemplate
from lollms.personality import APScript, AIPersonality
from lollms.prompting import LollmsContextDetails
import subprocess
if not PackageManager.check_package_installed("pptx"):
    PackageManager.install_package("pptx")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import Callable, Any
# Helper functions
class PowerPointBuilder:
    """
    A class for building PowerPoint slides programmatically using the python-pptx library.

    Attributes:
        presentation (Presentation): The PowerPoint presentation object.

    Methods:
        add_slide(layout=0): Adds a new slide to the presentation.
        add_text(slide, text, left, top, width, height, font_size=18, font_name="Arial", bold=False, italic=False, color=RGBColor(0, 0, 0)): Adds a text box with formatted text to a slide.
        add_image(slide, image_path, left, top, width=None, height=None): Adds an image to a slide.
        set_background_color(slide, color): Sets the background color of a slide.
        add_shape(slide, shape_type, left, top, width, height): Adds a shape to a slide.
        add_transition(slide, transition_type): Adds a transition to a slide.
        save(file_name): Saves the presentation to a file.

    """

    def __init__(self):
        self.presentation = Presentation()

    def add_slide(self, layout=0):
        """
        Adds a new slide to the presentation.

        Args:
            layout (int): The index of the slide layout to use. Default is 0.

        Returns:
            Slide: The newly added slide object.

        """
        slide_layout = self.presentation.slide_layouts[layout]
        slide = self.presentation.slides.add_slide(slide_layout)
        return slide

    def add_text(self, slide, text, left, top, width, height, font_size=18, font_name="Arial", bold=False, italic=False, color=RGBColor(0, 0, 0)):
        """
        Adds a text box with formatted text to a slide.

        Args:
            slide (Slide): The slide to add the text box to.
            text (str): The text to add.
            left (float): The left position of the text box in inches.
            top (float): The top position of the text box in inches.
            width (float): The width of the text box in inches.
            height (float): The height of the text box in inches.
            font_size (int): The font size of the text. Default is 18.
            font_name (str): The font name of the text. Default is "Arial".
            bold (bool): Whether the text should be bold. Default is False.
            italic (bool): Whether the text should be italic. Default is False.
            color (RGBColor): The color of the text. Default is black (RGBColor(0, 0, 0)).

        Returns:
            TextBox: The newly added text box object.

        """
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text

        font = run.font
        font.size = Pt(font_size)
        font.name = font_name
        font.bold = bold
        font.italic = italic
        font.color.rgb = color

        return textbox

    def add_image(self, slide, image_path, left, top, width=None, height=None):
        """
        Adds an image to a slide.

        Args:
            slide (Slide): The slide to add the image to.
            image_path (str): The path to the image file.
            left (float): The left position of the image in inches.
            top (float): The top position of the image in inches.
            width (float): The width of the image in inches. Default is None (original width).
            height (float): The height of the image in inches. Default is None (original height).

        """
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width) if width else None, height=Inches(height) if height else None)

    def set_background_color(self, slide, color):
        """
        Sets the background color of a slide.

        Args:
            slide (Slide): The slide to set the background color of.
            color (tuple): The RGB color values as a tuple (e.g., (255, 255, 255) for white).

        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])

    def add_shape(self, slide, shape_type, left, top, width, height):
        """
        Adds a shape to a slide.

        Args:
            slide (Slide): The slide to add the shape to.
            shape_type (MSO_SHAPE): The type of shape to add (e.g., MSO_SHAPE.RECTANGLE).
            left (float): The left position of the shape in inches.
            top (float): The top position of the shape in inches.
            width (float): The width of the shape in inches.
            height (float): The height of the shape in inches.

        Returns:
            Shape: The newly added shape object.

        """
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        return shape

    def add_transition(self, slide, transition_type):
        """
        Adds a transition to a slide.

        Args:
            slide (Slide): The slide to add the transition to.
            transition_type (str): The type of transition to add (e.g., "fade", "wipe").

        """
        # Note: Transitions are limited in python-pptx and may not be as flexible as PowerPoint's native transitions
        slide.slide_show_transition.type = transition_type

    def save(self, file_name):
        """
        Saves the presentation to a file.

        Args:
            file_name (str): The name of the file to save the presentation to.

        """
        self.presentation.save(file_name)
class Processor(APScript):
    """
    A class that processes model inputs and outputs.

    Inherits from APScript.
    """
    def __init__(
                 self, 
                 personality: AIPersonality,
                 callback = None,
                ) -> None:
        self.memory = []
        self.callback = None
        # Example entries
        #       {"name":"make_scripted","type":"bool","value":False, "help":"Makes a scriptred AI that can perform operations using python script"},
        #       {"name":"make_scripted","type":"bool","value":False, "help":"Makes a scriptred AI that can perform operations using python script"},
        # Supported types:
        # str, int, float, bool, list
        # options can be added using : "options":["option1","option2"...]        
        personality_config_template = ConfigTemplate(
            [
            ]
            )
        personality_config_vals = BaseConfig.from_template(personality_config_template)

        personality_config = TypedConfig(
            personality_config_template,
            personality_config_vals
        )
        super().__init__(
                            personality,
                            personality_config,
                            [
                                {
                                    "name": "idle",
                                    "commands": { # list of commands (don't forget to add these to your config.yaml file)
                                        "help":self.help,
                                    },
                                    "default": None
                                },                           
                            ],
                            callback=callback
                        )
        
    def install(self):
        super().install()
        
        # requirements_file = self.personality.personality_package_path / "requirements.txt"
        # Install dependencies using pip from requirements.txt
        # subprocess.run(["pip", "install", "--upgrade", "-r", str(requirements_file)])      
        ASCIIColors.success("Installed successfully")        

    def help(self, prompt="", full_context=""):
        self.set_message_content(self.personality.help)
    
    def add_file(self, path, client, callback=None):
        """
        Here we implement the file reception handling
        """
        super().add_file(path, client, callback)

    def get_description(self):
        return """
# PowerPointBuilder

A class for building PowerPoint slides programmatically using the python-pptx library.

## Attributes

- `presentation (Presentation)`: The PowerPoint presentation object.

## Methods

- `add_slide(layout=0)`: Adds a new slide to the presentation.
- `add_text(slide, text, left, top, width, height, font_size=18, font_name="Arial", bold=False, italic=False, color=RGBColor(0, 0, 0))`: Adds a text box with formatted text to a slide.
- `add_image(slide, image_path, left, top, width=None, height=None)`: Adds an image to a slide.
- `set_background_color(slide, color)`: Sets the background color of a slide.
- `add_shape(slide, shape_type, left, top, width, height)`: Adds a shape to a slide.
- `add_transition(slide, transition_type)`: Adds a transition to a slide.
- `save(file_name)`: Saves the presentation to a file.
"""

    from lollms.client_session import Client
    def run_workflow(self,  context_details:LollmsContextDetails=None, client:Client=None,  callback: Callable[[str | list | None, MSG_OPERATION_TYPE, str, AIPersonality| None], bool]=None):
        """
        This function generates code based on the given parameters.

        Args:
            context_details (dict): A dictionary containing the following context details for code generation:
                - conditionning (str): The conditioning information.
                - documentation (str): The documentation information.
                - knowledge (str): The knowledge information.
                - user_description (str): The user description information.
                - discussion_messages (str): The discussion messages information.
                - positive_boost (str): The positive boost information.
                - negative_boost (str): The negative boost information.
                - current_language (str): The force language information.
                - fun_mode (str): The fun mode conditionning text
                - ai_prefix (str): The AI prefix information.
            client_id: The client ID for code generation.
            callback (function, optional): The callback function for code generation.

        Returns:
            None
        """
        prompt = context_details.prompt
        previous_discussion_text = context_details.discussion_messages

        self.callback = callback
        self.personality.info("Generating")
        answer = self.multichoice_question("classify the user message",[
                                                "The user is providing information about the document to build",
                                                "The user is asking to start building the document",
                                                "The user is updating information about the document",
                                                "The user is asking for clarification",
                                                "The user is asking a question not related to the document"
                                            ],prompt)
        if answer==0:# providing information about the document to build
            # ask the ai to reformulate the promptelif answer==1: # ask the ai to start building the document
            self.memory.append(self.fastgen(f"Reformulate the user prompt in form of a list of entries that describe the request.\nprompt:{prompt}"))
            self.set_message_content("Information assimilated and stored to the memory. Do you want to add more information or do you want me to start generating the document?")            
        elif answer==1:
            if self.multichoice_question("classify the prompt",[
                                                "The prompt is do not contain useful information for the generation",
                                                "The prompt contains useful information for the generation",
                                            ],prompt):
                self.memory.append(self.fastgen(f"Reformulate the user prompt in form of a list of entries that describe the request.\nprompt:{prompt}"))

            
        else:
            out = self.fast_gen(previous_discussion_text)
            self.set_message_content(out)
        return out

